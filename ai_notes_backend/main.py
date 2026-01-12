import json
from typing import List
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
import os
import re
from pptx import Presentation
from pydantic import BaseModel
from transformers import pipeline
import pdfplumber


app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)



class MCQ(BaseModel):
    question: str
    options: List[str]
    correctIndex: int

summarizer = pipeline(
    "summarization",
    model="facebook/bart-large-cnn"
)

text_generator = pipeline(
    "text2text-generation",
    model="google/flan-t5-base"
)


UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

def clean_text(text: str) -> str:
    text = re.sub(r'\n+', '\n', text)
    text = re.sub(r'\s+', ' ', text)
    text = text.replace('\x00', '')
    return text.strip()

def chunk_text(text: str, max_words: int = 350):
    words = text.split()
    return [
        " ".join(words[i:i + max_words])
        for i in range(0, len(words), max_words)
    ]


def extract_pdf_text(path: str) -> str:
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text(x_tolerance=2)
            if page_text:
                text += page_text + "\n"
    return clean_text(text)

def extract_ppt_text(path: str) -> str:
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return clean_text(text)

# ------------------ AI LOGIC ------------------
def generate_summary(text: str) -> str:
    chunks = chunk_text(text)
    summaries = []

    for chunk in chunks:
        result = summarizer(
            chunk,
            max_length=130,
            min_length=60,
            do_sample=False
        )[0]["summary_text"]
        summaries.append(result)

    return " ".join(summaries)

def generate_notes(text: str) -> str:
    prompt = f"""
    You are an expert academic tutor.

    Convert the following content into:
    - Clear and concise study notes
    - Bullet points
    - Key concepts only
    - Easy for exams and revision

    Content:
    {text}
    """

    result = text_generator(
        prompt,
        max_length=512,
        do_sample=False
    )[0]["generated_text"]

    return result

def generate_quiz(text: str):
    prompt = f"""
Create 5 multiple choice questions.

Format:
Q: question
A) option
B) option
C) option
D) option
ANSWER: A

Text:
{text[:1200]}
"""

    output = text_generator(
        prompt,
        max_new_tokens=400,
        do_sample=False
    )[0]["generated_text"]

    print("\n===== RAW QUIZ OUTPUT =====\n")
    print(output)
    print("\n===========================\n")

    return parse_mcq_text(output)



def parse_mcq_text(raw: str):
    questions = []

    blocks = re.split(r"\n\s*Q[:\-]\s*", raw, flags=re.IGNORECASE)

    for block in blocks:
        if not block.strip():
            continue

        lines = [l.strip() for l in block.splitlines() if l.strip()]

        try:
            question = lines[0]

            option_lines = [l for l in lines if re.match(r"^[A-Da-d][\).]", l)]
            if len(option_lines) < 4:
                continue

            options = [
                re.sub(r"^[A-Da-d][\).]\s*", "", opt)
                for opt in option_lines[:4]
            ]

            answer_line = next(
                (l for l in lines if "answer" in l.lower()), None
            )
            if not answer_line:
                continue

            letter = re.search(r"[A-Da-d]", answer_line).group().upper()
            correct_index = ord(letter) - ord("A")

            questions.append({
                "question": question,
                "options": options,
                "correctIndex": correct_index
            })

        except Exception as e:
            print("MCQ PARSE ERROR:", e)

    return questions


    # ---------- SAFE JSON EXTRACTION ----------
    try:
        json_text = re.search(r"\[.*\]", result, re.S).group()
        quiz = json.loads(json_text)
        return quiz
    except Exception as e:
        print("Quiz parsing error:", e)
        return []

# ------------------ API ENDPOINTS ------------------
@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    path = f"{UPLOAD_DIR}/{file.filename}"
    with open(path, "wb") as f:
        f.write(await file.read())

    return {
        "filename": file.filename,
        "path": path
    }

@app.post("/process")
async def process_file(file: UploadFile = File(...)):
    path = f"{UPLOAD_DIR}/{file.filename}"
    with open(path, "wb") as f:
        f.write(await file.read())

    if file.filename.lower().endswith(".pdf"):
        raw_text = extract_pdf_text(path)
    else:
        raw_text = extract_ppt_text(path)

    summary = generate_summary(raw_text)
    notes = generate_notes(summary)
    quiz = generate_quiz(summary)

    print("QUIZ COUNT:", len(quiz))  # 🔍 DEBUG

    return {
        "summary": summary,
        "notes": notes,
        "quiz": quiz
    }




# py -m uvicorn main:app --reload